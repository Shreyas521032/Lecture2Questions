import streamlit as st
import google.generativeai as genai
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from datetime import datetime
from io import BytesIO
import json
import random

# ===== Configuration =====
st.set_page_config(
    page_title="EduGenius - AI Question Generator",
    page_icon="🧠",
    layout="wide",
    initial_sidebar_state="expanded"
)

# ===== Custom CSS =====
st.markdown("""
<style>
    /* Main Styling */
    :root {
        --primary: #4361ee;
        --primary-light: #4895ef;
        --secondary: #3a0ca3;
        --accent: #f72585;
        --text-dark: #2b2d42;
        --text-light: #8d99ae;
        --bg-light: #f8f9fa;
        --bg-medium: #e9ecef;
        --success: #4cc9a7;
        --warning: #ffc107;
        --error: #ef476f;
        --border-radius: 8px;
    }
    
    .main-title {
        font-size: 2.2rem;
        font-weight: 800;
        color: var(--primary);
        margin-bottom: 0.5rem;
        padding-bottom: 0.5rem;
        border-bottom: 2px solid var(--primary-light);
    }
    
    .subtitle {
        font-size: 1.1rem;
        color: var(--text-light);
        margin-bottom: 2rem;
    }
    
    .card {
        background-color: white;
        border-radius: var(--border-radius);
        padding: 1.5rem;
        box-shadow: 0 4px 6px rgba(0,0,0,0.05);
        margin-bottom: 1.5rem;
        border-top: 4px solid var(--primary);
    }
    
    .card-title {
        font-size: 1.4rem;
        font-weight: 600;
        color: var(--primary);
        margin-bottom: 1rem;
        display: flex;
        align-items: center;
        gap: 0.5rem;
    }
    
    .card-title svg {
        width: 1.4rem;
        height: 1.4rem;
    }
    
    /* Question Styling */
    .question-container {
        background-color: var(--bg-light);
        border-radius: var(--border-radius);
        padding: 1.5rem;
        margin-bottom: 1.5rem;
        border-left: 4px solid var(--primary);
        box-shadow: 0 2px 4px rgba(0,0,0,0.03);
    }
    
    .question-header {
        display: flex;
        justify-content: space-between;
        margin-bottom: 0.75rem;
        border-bottom: 1px solid var(--bg-medium);
        padding-bottom: 0.5rem;
    }
    
    .question-number {
        font-weight: 700;
        color: var(--primary);
        font-size: 1.1rem;
    }
    
    .question-difficulty {
        padding: 0.2rem 0.6rem;
        border-radius: 12px;
        font-size: 0.8rem;
        font-weight: 600;
    }
    
    .difficulty-easy {
        background-color: #d7f9e9;
        color: #0d8a53;
    }
    
    .difficulty-medium {
        background-color: #fff8dd;
        color: #b88400;
    }
    
    .difficulty-hard {
        background-color: #ffe5e5;
        color: #c11c1c;
    }
    
    .question-text {
        font-size: 1.1rem;
        margin-bottom: 1rem;
        color: var(--text-dark);
        line-height: 1.5;
    }
    
    .options {
        display: grid;
        grid-template-columns: 1fr 1fr;
        gap: 0.75rem;
        margin: 1rem 0;
    }
    
    .option {
        padding: 0.75rem;
        background-color: white;
        border: 1px solid var(--bg-medium);
        border-radius: var(--border-radius);
        cursor: pointer;
        transition: all 0.2s;
    }
    
    .option:hover {
        border-color: var(--primary-light);
        box-shadow: 0 2px 4px rgba(67, 97, 238, 0.1);
    }
    
    .option-correct {
        border-color: var(--success);
        background-color: #f0fcf8;
    }
    
    .answer-section {
        margin-top: 1rem;
        border-top: 1px dashed var(--bg-medium);
        padding-top: 1rem;
    }
    
    .answer-label {
        font-weight: 600;
        color: var(--success);
        margin-bottom: 0.5rem;
    }
    
    .explanation {
        background-color: #f8f9fa;
        padding: 0.75rem;
        border-radius: var(--border-radius);
        border-left: 3px solid var(--text-light);
        font-size: 0.95rem;
        margin-top: 0.75rem;
    }
    
    /* Stats Cards */
    .stat-container {
        display: flex;
        gap: 1rem;
        margin-bottom: 1.5rem;
    }
    
    .stat-card {
        flex: 1;
        background: white;
        border-radius: var(--border-radius);
        padding: 1rem;
        box-shadow: 0 2px 4px rgba(0,0,0,0.05);
        display: flex;
        flex-direction: column;
        align-items: center;
        text-align: center;
        border-bottom: 3px solid var(--primary);
    }
    
    .stat-title {
        color: var(--text-light);
        font-size: 0.9rem;
        margin-bottom: 0.5rem;
    }
    
    .stat-value {
        font-size: 2rem;
        font-weight: 700;
        color: var(--primary);
    }
    
    .stat-subtitle {
        font-size: 0.8rem;
        color: var(--text-light);
    }
    
    /* Buttons and Inputs */
    .custom-button {
        background-color: var(--primary);
        color: white;
        border: none;
        padding: 0.6rem 1.2rem;
        border-radius: var(--border-radius);
        font-weight: 600;
        cursor: pointer;
        transition: all 0.3s;
        display: inline-flex;
        align-items: center;
        gap: 0.5rem;
    }
    
    .custom-button:hover {
        background-color: var(--secondary);
        transform: translateY(-2px);
        box-shadow: 0 4px 8px rgba(0,0,0,0.1);
    }
    
    .custom-button-secondary {
        background-color: white;
        color: var(--primary);
        border: 1px solid var(--primary);
    }
    
    .custom-button-secondary:hover {
        background-color: var(--bg-light);
        color: var(--secondary);
        border-color: var(--secondary);
    }
    
    /* Tags and Badges */
    .tag {
        display: inline-block;
        padding: 0.3rem 0.6rem;
        border-radius: 50px;
        font-size: 0.8rem;
        font-weight: 600;
        margin-right: 0.5rem;
        margin-bottom: 0.5rem;
    }
    
    .tag-mcq {
        background-color: #e7f5ff;
        color: #1971c2;
    }
    
    .tag-short {
        background-color: #fff9db;
        color: #e67700;
    }
    
    .tag-fill {
        background-color: #f3f0ff;
        color: #5f3dc4;
    }
    
    /* Analytics */
    .analytics-container {
        background-color: white;
        border-radius: var(--border-radius);
        padding: 1.5rem;
        box-shadow: 0 4px 6px rgba(0,0,0,0.05);
        margin-bottom: 1.5rem;
    }
    
    .analytics-title {
        font-size: 1.2rem;
        font-weight: 600;
        color: var(--text-dark);
        margin-bottom: 1rem;
        padding-bottom: 0.5rem;
        border-bottom: 1px solid var(--bg-medium);
    }
    
    /* Streamlit Overrides */
    div.stButton > button {
        width: 100%;
        background-color: var(--primary);
        color: white;
        font-weight: 600;
        border-radius: var(--border-radius);
        border: none;
        padding: 0.5rem 1rem;
        transition: all 0.3s;
    }
    
    div.stButton > button:hover {
        background-color: var(--secondary);
        border: none;
    }
    
    .stTabs [data-baseweb="tab-list"] {
        gap: 8px;
    }
    
    .stTabs [data-baseweb="tab"] {
        height: 50px;
        border-radius: var(--border-radius);
        background-color: white;
        padding: 0 16px;
        box-shadow: 0 2px 4px rgba(0,0,0,0.05);
    }
    
    .stTabs [aria-selected="true"] {
        background-color: var(--primary);
        color: white;
    }
    
    /* Report Styling */
    .report-container {
        background-color: white;
        border-radius: var(--border-radius);
        padding: 20px;
        margin-bottom: 20px;
        box-shadow: 0 2px 4px rgba(0,0,0,0.05);
    }
    
    .report-header {
        display: flex;
        justify-content: space-between;
        align-items: center;
        margin-bottom: 15px;
        padding-bottom: 10px;
        border-bottom: 1px solid var(--bg-medium);
    }
    
    .report-title {
        font-size: 1.2rem;
        font-weight: 600;
        color: var(--primary);
    }
    
    /* Datatables */
    .dataframe {
        border-collapse: collapse;
        width: 100%;
        border: none;
    }
    
    .dataframe th {
        background-color: var(--primary);
        color: white;
        font-weight: 600;
        text-align: left;
        padding: 12px;
    }
    
    .dataframe td {
        padding: 10px 12px;
        border-bottom: 1px solid var(--bg-medium);
    }
    
    .dataframe tr:hover {
        background-color: var(--bg-light);
    }
</style>
""", unsafe_allow_html=True)

# ===== Session State Initialization =====
if 'generation_history' not in st.session_state:
    st.session_state.generation_history = []

if 'user_settings' not in st.session_state:
    st.session_state.user_settings = {
        'api_key': '',
        'model': 'gemini-1.5-pro-latest',
        'theme': 'blue',
        'show_explanations': True
    }

if 'question_bank' not in st.session_state:
    st.session_state.question_bank = []

# ===== Helper Functions =====
def extract_text(file):
    try:
        if file.type == "text/plain":
            return file.getvalue().decode("utf-8")
        
        elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            from docx import Document
            doc = Document(BytesIO(file.read()))
            return "\n".join([para.text for para in doc.paragraphs])
        
        elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            from pptx import Presentation
            prs = Presentation(BytesIO(file.read()))
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
        
        elif file.type == "application/pdf":
            import PyPDF2
            pdf_reader = PyPDF2.PdfReader(BytesIO(file.read()))
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
            return text
        
        else:
            st.error("Unsupported file type")
            return ""
    
    except Exception as e:
        st.error(f"Error processing file: {str(e)}")
        return ""

def generate_with_gemini(text, question_type, difficulty, api_key, num_questions=5, topics=None):
    try:
        genai.configure(api_key=api_key)
        
        # Use the selected model
        model_name = st.session_state.user_settings.get('model', 'gemini-1.5-pro-latest')
        model = genai.GenerativeModel(model_name)
        
        # Build topic-specific prompt if topics are provided
        topic_prompt = ""
        if topics and len(topics) > 0:
            topic_prompt = f"Focus on these specific topics: {', '.join(topics)}.\n"
            
        
        response = model.generate_content(prompt)
        
        if not response.text:
            raise ValueError("Empty response from Gemini API")
        
        # Extract JSON from response
        import re
        json_match = re.search(r'```json\s*([\s\S]*?)\s*```', response.text)
        
        if json_match:
            json_str = json_match.group(1)
            try:
                questions_data = json.loads(json_str)
                return questions_data
            except json.JSONDecodeError:
                # Fallback to text response if JSON parsing fails
                return response.text
        else:
            # If no JSON format is found, return the text as is
            return response.text
    
    except Exception as e:
        st.error(f"Gemini API error: {str(e)}")
        return None

def parse_questions(raw_questions):
    """Parse questions from raw text if JSON format failed"""
    if isinstance(raw_questions, list):
        # Already in proper format
        return raw_questions
    
    # Try to parse manually
    import re
    questions = []
    
    # Simple pattern matching for question components
    q_blocks = re.split(r'\n\s*\d+\.\s+', raw_questions)
    if len(q_blocks) <= 1:
        return []
    
    for i, block in enumerate(q_blocks[1:], 1):  # Skip first empty element
        q_data = {
            "id": i,
            "type": "Unknown",
            "difficulty": "medium"
        }
        
        # Extract question
        q_match = re.search(r'(.*?)(?:Options:|Correct Answer:|$)', block, re.DOTALL)
        if q_match:
            q_data["question"] = q_match.group(1).strip()
            
        # Extract options
        opt_match = re.search(r'Options:(.*?)(?:Correct Answer:|$)', block, re.DOTALL)
        if opt_match:
            options_text = opt_match.group(1).strip()
            options = re.findall(r'[A-D]\)\s*(.*?)(?=\s*[A-D]\)|$)', options_text + " ")
            q_data["options"] = [opt.strip() for opt in options]
            q_data["type"] = "MCQ"
        else:
            q_data["type"] = "Short Answer"
        
        # Extract correct answer
        ans_match = re.search(r'Correct Answer:(.*?)(?:Explanation:|$)', block, re.DOTALL)
        if ans_match:
            q_data["correctAnswer"] = ans_match.group(1).strip()
            
        # Extract explanation
        exp_match = re.search(r'Explanation:(.*?)$', block, re.DOTALL)
        if exp_match:
            q_data["explanation"] = exp_match.group(1).strip()
            
        questions.append(q_data)
    
    return questions

def get_topic_statistics(questions):
    """Extract topic statistics from questions"""
    topics = {}
    
    for q in questions:
        if "topics" in q and q["topics"]:
            for topic in q["topics"]:
                if topic in topics:
                    topics[topic] += 1
                else:
                    topics[topic] = 1
    
    return topics

def generate_question_distribution(history):
    """Generate data for question distribution chart"""
    question_types = {}
    difficulties = {}
    
    for entry in history:
        qtype = entry.get('type', 'Unknown')
        if qtype in question_types:
            question_types[qtype] += 1
        else:
            question_types[qtype] = 1
            
        diff = entry.get('difficulty', 'Medium')
        if diff in difficulties:
            difficulties[diff] += 1
        else:
            difficulties[diff] = 1
    
    return question_types, difficulties

def generate_pdf_report(questions, metadata):
    """Generate a PDF report of questions"""
    from reportlab.lib.pagesizes import A4
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib import colors
    from io import BytesIO
    
    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4, rightMargin=72, leftMargin=72, topMargin=72, bottomMargin=72)
    
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(name='Heading1', fontSize=16, spaceAfter=12, textColor=colors.HexColor("#4361ee")))
    styles.add(ParagraphStyle(name='Question', fontSize=12, spaceAfter=6, fontName='Helvetica-Bold'))
    styles.add(ParagraphStyle(name='Option', fontSize=11, leftIndent=20))
    styles.add(ParagraphStyle(name='Answer', fontSize=11, textColor=colors.HexColor("#4cc9a7"), fontName='Helvetica-Bold'))
    styles.add(ParagraphStyle(name='Explanation', fontSize=10, leftIndent=20, backgroundColor=colors.HexColor("#f8f9fa")))
    
    # Build document content
    content = []
    
    # Add title
    title = Paragraph(f"Question Set: {metadata['file']}", styles['Heading1'])
    content.append(title)
    
    # Add metadata
    meta_data = [
        [Paragraph("Date:", styles['BodyText']), Paragraph(metadata['timestamp'], styles['BodyText'])],
        [Paragraph("Question Type:", styles['BodyText']), Paragraph(metadata['type'], styles['BodyText'])],
        [Paragraph("Difficulty:", styles['BodyText']), Paragraph(metadata['difficulty'], styles['BodyText'])],
        [Paragraph("Number of Questions:", styles['BodyText']), Paragraph(str(metadata['count']), styles['BodyText'])]
    ]
    meta_table = Table(meta_data, colWidths=[100, 300])
    meta_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (0, -1), colors.HexColor("#f8f9fa")),
        ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor("#e9ecef"))
    ]))
    content.append(meta_table)
    content.append(Spacer(1, 12))
    
    # Add questions
    for i, q in enumerate(questions, 1):
        # Question text
        q_text = Paragraph(f"{i}. {q['question']}", styles['Question'])
        content.append(q_text)
        content.append(Spacer(1, 6))
        
        # Options for MCQ
        if 'options' in q and q['options']:
            for j, opt in enumerate(q['options']):
                option_letter = chr(65 + j)  # A, B, C, D...
                opt_text = Paragraph(f"{option_letter}) {opt}", styles['Option'])
                content.append(opt_text)
        
        # Answer
        if 'correctAnswer' in q and q['correctAnswer']:
            ans_text = Paragraph(f"Correct Answer: {q['correctAnswer']}", styles['Answer'])
            content.append(ans_text)
            
        # Explanation
        if 'explanation' in q and q['explanation']:
            exp_text = Paragraph(f"Explanation: {q['explanation']}", styles['Explanation'])
            content.append(exp_text)
            
        content.append(Spacer(1, 12))
    
    # Build PDF
    doc.build(content)
    buffer.seek(0)
    return buffer

def display_question(q, show_explanation=True):
    """Display a formatted question"""
    q_type = q.get('type', 'Unknown')
    difficulty = q.get('difficulty', 'medium')
    
    # Difficulty class
    difficulty_class = "difficulty-medium"
    if difficulty.lower() == "easy":
        difficulty_class = "difficulty-easy"
    elif difficulty.lower() == "hard":
        difficulty_class = "difficulty-hard"
        
    # Question header
    st.markdown(f"""
    <div class="question-header">
        <div class="question-number">Question {q.get('id', '')}</div>
        <div class="question-difficulty {difficulty_class}">{difficulty.title()}</div>
    </div>
    """, unsafe_allow_html=True)
    
    # Question text
    st.markdown(f'<div class="question-text">{q.get("question", "")}</div>', unsafe_allow_html=True)
    
    # Options for MCQ
    if q_type.upper() == "MCQ" and 'options' in q and q['options']:
        option_cols = st.columns(2)
        for i, option in enumerate(q['options']):
            col_idx = i % 2
            with option_cols[col_idx]:
                option_letter = chr(65 + i)  # A, B, C, D...
                option_class = "option"
    
    # Display answer section if enabled
    if show_explanation:
        st.markdown('<div class="answer-section">', unsafe_allow_html=True)
        st.markdown(f'<div class="answer-label">Correct Answer:</div>', unsafe_allow_html=True)
        st.markdown(f'<div>{q.get("correctAnswer", "")}</div>', unsafe_allow_html=True)
        
        if 'explanation' in q and q['explanation']:
            st.markdown(f'<div class="explanation">{q.get("explanation", "")}</div>', unsafe_allow_html=True)
        st.markdown('</div>', unsafe_allow_html=True)
        
    # Display topics if available
    if 'topics' in q and q['topics']:
        st.markdown("<p>Topics:</p>", unsafe_allow_html=True)
        topic_html = ""
        for topic in q['topics']:
            topic_html += f'<span class="tag tag-mcq">{topic}</span>'
        st.markdown(topic_html, unsafe_allow_html=True)

# ===== Sidebar Navigation =====
with st.sidebar:
    st.image("https://i.imgur.com/bGsBgPA.png", width=50)
    st.markdown("### EduGenius")
    st.markdown("AI-Powered Question Generator")
    
    st.markdown("---")
    
    tab = st.radio(
        "Navigation",
        ["🧠 Generate Questions", "📊 Analytics", "🏦 Question Bank", "⚙️ Settings"],
        captions=["Create new questions", "View insights", "Saved questions", "Configure app"]
    )
    
    st.markdown("---")
    
    # Settings section in sidebar
    if tab == "⚙️ Settings":
        st.markdown("### API Configuration")
        api_key = st.text_input(
            "Gemini API Key", 
            value=st.session_state.user_settings['api_key'],
            type="password", 
            help="Get your key from Google AI Studio"
        )
        
        model = st.selectbox(
            "Model", 
            ["gemini-1.5-pro-latest", "gemini-1.0-pro"],
            index=0 if st.session_state.user_settings['model'] == "gemini-1.5-pro-latest" else 1
        )
        
        show_explanations = st.checkbox(
            "Show explanations by default", 
            value=st.session_state.user_settings['show_explanations']
        )
        
        if st.button("Save Settings"):
            st.session_state.user_settings['api_key'] = api_key
            st.session_state.user_settings['model'] = model
            st.session_state.user_settings['show_explanations'] = show_explanations
            st.success("Settings saved!")
    
    # Quick stats in sidebar (except on analytics page)
    if tab != "📊 Analytics":
        st.markdown("### Quick Stats")
        
        total_gens = len(st.session_state.generation_history)
        total_questions = sum(entry.get('count', 0) for entry in st.session_state.generation_history)
        
        st.markdown(f"""
        <div style="display: flex; justify-content: space-between; margin-bottom: 10px;">
            <div style="font-weight: 600; color: #4361ee;">Total Generations:</div>
            <div>{total_gens}</div>
        </div>
        <div style="display: flex; justify-content: space-between; margin-bottom: 10px;">
            <div style="font-weight: 600; color: #4361ee;">Total Questions:</div>
            <div>{total_questions}</div>
        </div>
        """, unsafe_allow_html=True)

# ===== Main Content =====
if tab == "🧠 Generate Questions":
    st.markdown('<h1 class="main-title">🧠 AI Question Generator</h1>', unsafe_allow_html=True)
    st.markdown('<p class="subtitle">Generate customized questions from your learning materials</p>', unsafe_allow_html=True)
    
    col1, col2 = st.columns([1, 2])
    
    with col1:
        st.markdown('<div class="card">', unsafe_allow_html=True)
        st.markdown('<div class="card-title">📄 Upload Material</div>', unsafe_allow_html=True)
        
        uploaded_file = st.file_uploader("Choose file", type=["txt", "docx", "pptx", "pdf"])
        
        if uploaded_file:
            st.success(f"File uploaded: {uploaded_file.name}")
            
            with st.expander("Configuration", expanded=True):
                question_type = st.selectbox(
                    "Question Type", 
                    ["MCQ", "Fill-in-the-Blank", "Short Answer"]
                )
                
                difficulty = st.selectbox("Difficulty", ["Easy", "Medium", "Hard"])
                
                num_questions = st.slider(
                    "Number of Questions", 
                    min_value=3, 
                    max_value=20, 
                    value=5
                )
                
                # New: Topic extraction
                st.markdown("##### Focus Topics (Optional)")
                topic_input = st.text_input(
                    "Enter topics separated by commas", 
                    help="Leave empty to generate questions on all topics"
                )
                topics = [t.strip() for t in topic_input.split(",")] if topic_input else []
            
            api_key = st.session_state.user_settings.get('api_key', '')
            if not api_key:
                st.warning("Please set your Gemini API key in Settings first")
            
            col1, col2 = st.columns(2)
            with col1:
                generate_btn = st.button("✨ Generate Questions", use_container_width=True)
            with col2:
                clear_btn = st.button("🗑️ Clear", use_container_width=True)
            
            if clear_btn:
                if 'questions' in st.session_state:
                    del st.session_state.questions
            
            if generate_btn:
                if not api_key:
                    st.error("Please enter your Gemini API key in Settings")
                else:
                    with st.spinner("Generating questions..."):
                        text = extract_text(uploaded_file)
                        if text:
                            questions_data = generate_with_gemini(
                                text, 
                                question_type, 
                                difficulty, 
                                api_key, 
                                num_questions=num_questions,
                                topics=topics
                            )
                            
                            # Parse questions if not in proper format
                            if questions_data:
                                if isinstance(questions_data, str):
                                    st.session_state.questions = parse_questions(questions_data)
                                else:
                                    st.session_state.questions = questions_data
                                
                                # Store in generation history
                                timestamp = datetime.now().strftime("%Y-%m-%d %H:%M")
                                history_entry = {
                                    'timestamp': timestamp,
                                    'file': uploaded_file.name,
                                    'type': question_type,
                                    'difficulty': difficulty,
                                    'count': len(st.session_state.questions),
                                    'questions': st.session_state.questions
                                }
                                st.session_state.generation_history.append(history_entry)
                                
                                # Add to question bank
                                for q in st.session_state.questions:
                                    # Add a unique identifier and source info
                                    q['source'] = uploaded_file.name
                                    q['created'] = timestamp
                                    st.session_state.question_bank.append(q)
                                
                                st.success(f"Generated {len(st.session_state.questions)} questions!")
                            else:
                                st.error("Failed to generate questions. Please try again.")
                        else:
                            st.error("Could not extract text from the file")
        
        st.markdown('</div>', unsafe_allow_html=True)
    
    with col2:
        if 'questions' in st.session_state and st.session_state.questions:
            st.markdown('<div class="card">', unsafe_allow_html=True)
            st.markdown('<div class="card-title">📝 Generated Questions</div>', unsafe_allow_html=True)
            
            # Export options
            col1, col2, col3 = st.columns([1, 1, 1])
            with col1:
                show_explanations = st.toggle(
                    "Show Explanations", 
                    value=st.session_state.user_settings.get('show_explanations', True)
                )
            
            with col3:
                if st.button("📥 Export PDF"):
                    # Get the most recent generation
                    latest = st.session_state.generation_history[-1]
                    
                    # Generate PDF
                    pdf_buffer = generate_pdf_report(
                        latest['questions'],
                        {
                            'file': latest['file'],
                            'timestamp': latest['timestamp'],
                            'type': latest['type'],
                            'difficulty': latest['difficulty'],
                            'count': latest['count']
                        }
                    )
                    
                    # Offer download
                    st.download_button(
                        label="Download PDF",
                        data=pdf_buffer,
                        file_name=f"questions_{datetime.now().strftime('%Y%m%d_%H%M')}.pdf",
                        mime="application/pdf"
                    )
            
            # Display questions
            st.markdown("<div class='question-list'>", unsafe_allow_html=True)
            for q in st.session_state.questions:
                st.markdown("<div class='question-container'>", unsafe_allow_html=True)
                display_question(q, show_explanation=show_explanations)
                st.markdown("</div>", unsafe_allow_html=True)
            st.markdown("</div>", unsafe_allow_html=True)
            
            st.markdown('</div>', unsafe_allow_html=True)
        else:
            st.markdown("""
            <div style="
                height: 300px; 
                display: flex; 
                flex-direction: column; 
                justify-content: center; 
                align-items: center; 
                text-align: center;
                background-color: #f8f9fa;
                border-radius: 8px;
                margin-top: 20px;">
                <img src="https://i.imgur.com/YSjyLmm.png" width="120">
                <h3 style="margin-top: 20px; color: #4361ee;">Upload material and generate questions</h3>
                <p style="color: #8d99ae; max-width: 400px;">
                    Upload your learning material on the left panel and click Generate to create AI-powered questions
                </p>
            </div>
            """, unsafe_allow_html=True)

elif tab == "📊 Analytics":
    st.markdown('<h1 class="main-title">📊 Generation Analytics</h1>', unsafe_allow_html=True)
    st.markdown('<p class="subtitle">Insights and statistics about your generated questions</p>', unsafe_allow_html=True)
    
    if len(st.session_state.generation_history) > 0:
        # Stats Cards
        col1, col2, col3, col4 = st.columns(4)
        
        with col1:
            st.markdown("""
            <div class="stat-card">
                <div class="stat-title">Total Generations</div>
                <div class="stat-value">{}</div>
                <div class="stat-subtitle">Question sets</div>
            </div>
            """.format(len(st.session_state.generation_history)), unsafe_allow_html=True)
        
        with col2:
            total_questions = sum(entry.get('count', 0) for entry in st.session_state.generation_history)
            st.markdown("""
            <div class="stat-card">
                <div class="stat-title">Total Questions</div>
                <div class="stat-value">{}</div>
                <div class="stat-subtitle">Generated</div>
            </div>
            """.format(total_questions), unsafe_allow_html=True)
        
        with col3:
            # Calculate average questions per generation
            avg_questions = total_questions / len(st.session_state.generation_history) if len(st.session_state.generation_history) > 0 else 0
            st.markdown("""
            <div class="stat-card">
                <div class="stat-title">Average Questions</div>
                <div class="stat-value">{:.1f}</div>
                <div class="stat-subtitle">Per generation</div>
            </div>
            """.format(avg_questions), unsafe_allow_html=True)
        
        with col4:
            # Calculate topic diversity
            all_topics = []
            for entry in st.session_state.generation_history:
                for q in entry.get('questions', []):
                    if 'topics' in q:
                        all_topics.extend(q['topics'])
            unique_topics = len(set(all_topics))
            
            st.markdown("""
            <div class="stat-card">
                <div class="stat-title">Topic Diversity</div>
                <div class="stat-value">{}</div>
                <div class="stat-subtitle">Unique topics</div>
            </div>
            """.format(unique_topics), unsafe_allow_html=True)
        
        # Visualization Section
        st.markdown('<div class="card">', unsafe_allow_html=True)
        st.markdown('<div class="card-title">📈 Question Distribution</div>', unsafe_allow_html=True)
        
        # Create a flat list of all questions
        all_questions = []
        for entry in st.session_state.generation_history:
            for q in entry.get('questions', []):
                all_questions.append(q)
        
        # Get distribution data
        question_types, difficulties = generate_question_distribution(all_questions)
        
        # Create visualizations
        col1, col2 = st.columns(2)
        
        with col1:
            fig1 = px.pie(
                names=list(question_types.keys()),
                values=list(question_types.values()),
                title="Question Types",
                color_discrete_sequence=px.colors.sequential.Blues_r,
                hole=0.4
            )
            fig1.update_traces(textposition='inside', textinfo='percent+label')
            fig1.update_layout(
                legend=dict(orientation="h", yanchor="bottom", y=-0.3, xanchor="center", x=0.5),
                margin=dict(l=20, r=20, t=40, b=20)
            )
            st.plotly_chart(fig1, use_container_width=True)
        
        with col2:
            fig2 = px.bar(
                x=list(difficulties.keys()),
                y=list(difficulties.values()),
                title="Difficulty Distribution",
                color=list(difficulties.keys()),
                color_discrete_map={
                    'easy': '#4cc9a7',
                    'medium': '#ffc107',
                    'hard': '#ef476f',
                    'Easy': '#4cc9a7',
                    'Medium': '#ffc107',
                    'Hard': '#ef476f'
                }
            )
            fig2.update_layout(
                xaxis_title="Difficulty",
                yaxis_title="Count",
                showlegend=False,
                margin=dict(l=20, r=20, t=40, b=20)
            )
            st.plotly_chart(fig2, use_container_width=True)
        
        st.markdown('</div>', unsafe_allow_html=True)
        
        # Topic Analysis
        st.markdown('<div class="card">', unsafe_allow_html=True)
        st.markdown('<div class="card-title">🔍 Topic Analysis</div>', unsafe_allow_html=True)
        
        # Extract topics
        topics_analysis = {}
        for q in all_questions:
            if 'topics' in q and q['topics']:
                for topic in q['topics']:
                    if topic in topics_analysis:
                        topics_analysis[topic] += 1
                    else:
                        topics_analysis[topic] = 1
        
        # Sort topics by frequency
        sorted_topics = dict(sorted(topics_analysis.items(), key=lambda x: x[1], reverse=True))
        
        # Display top topics
        if sorted_topics:
            # Create a horizontal bar chart for topics
            topic_df = pd.DataFrame({
                'Topic': list(sorted_topics.keys())[:15],  # Top 15 topics
                'Count': list(sorted_topics.values())[:15]
            })
            
            fig = px.bar(
                topic_df, 
                y='Topic', 
                x='Count', 
                orientation='h',
                color='Count',
                color_continuous_scale=px.colors.sequential.Blues,
                title="Most Common Topics"
            )
            
            fig.update_layout(
                yaxis={'categoryorder':'total ascending'},
                margin=dict(l=20, r=20, t=40, b=20)
            )
            
            st.plotly_chart(fig, use_container_width=True)
        else:
            st.info("No topic data available for analysis")
        
        st.markdown('</div>', unsafe_allow_html=True)
        
        # Generation History
        st.markdown('<div class="card">', unsafe_allow_html=True)
        st.markdown('<div class="card-title">📜 Generation History</div>', unsafe_allow_html=True)
        
        # Create history dataframe
        history_data = []
        for entry in st.session_state.generation_history:
            history_data.append({
                'Date': entry.get('timestamp', ''),
                'File': entry.get('file', ''),
                'Type': entry.get('type', ''),
                'Difficulty': entry.get('difficulty', ''),
                'Questions': entry.get('count', 0)
            })
        
        history_df = pd.DataFrame(history_data).sort_values(by='Date', ascending=False)
        st.dataframe(history_df, use_container_width=True)
        st.markdown('</div>', unsafe_allow_html=True)
        
    else:
        st.info("No generation history found. Generate some questions first to see analytics.")

elif tab == "🏦 Question Bank":
    st.markdown('<h1 class="main-title">🏦 Question Bank</h1>', unsafe_allow_html=True)
    st.markdown('<p class="subtitle">Browse and manage your collection of generated questions</p>', unsafe_allow_html=True)
    
    if len(st.session_state.question_bank) > 0:
        # Filter options
        col1, col2, col3 = st.columns(3)
        
        with col1:
            # Get unique file sources
            sources = list(set([q.get('source', 'Unknown') for q in st.session_state.question_bank]))
            source_filter = st.selectbox("Filter by Source", ["All"] + sources)
        
        with col2:
            # Get question types
            types = list(set([q.get('type', 'Unknown') for q in st.session_state.question_bank]))
            type_filter = st.selectbox("Filter by Type", ["All"] + types)
        
        with col3:
            # Get difficulties
            difficulties = list(set([q.get('difficulty', 'medium') for q in st.session_state.question_bank]))
            difficulty_filter = st.selectbox("Filter by Difficulty", ["All"] + difficulties)
        
        # Apply filters
        filtered_questions = st.session_state.question_bank.copy()
        
        if source_filter != "All":
            filtered_questions = [q for q in filtered_questions if q.get('source', '') == source_filter]
        
        if type_filter != "All":
            filtered_questions = [q for q in filtered_questions if q.get('type', '') == type_filter]
        
        if difficulty_filter != "All":
            filtered_questions = [q for q in filtered_questions if q.get('difficulty', '') == difficulty_filter]
        
        # Search functionality
        search_query = st.text_input("🔍 Search Questions", help="Search by keywords in questions")
        if search_query:
            search_terms = search_query.lower().split()
            filtered_questions = [
                q for q in filtered_questions 
                if any(term in q.get('question', '').lower() for term in search_terms)
            ]
        
        # Display questions
        st.markdown(f"### Displaying {len(filtered_questions)} Questions")
        
        # Question display
        show_explanations = st.toggle("Show Explanations", value=st.session_state.user_settings.get('show_explanations', True))
        
        for q in filtered_questions:
            st.markdown("<div class='question-container'>", unsafe_allow_html=True)
            
            # Source info
            st.markdown(f"""
            <div style="margin-bottom: 10px; color: #8d99ae; font-size: 0.8rem;">
                Source: {q.get('source', 'Unknown')} | Created: {q.get('created', 'Unknown')}
            </div>
            """, unsafe_allow_html=True)
            
            display_question(q, show_explanation=show_explanations)
            st.markdown("</div>", unsafe_allow_html=True)
        
        if len(filtered_questions) == 0:
            st.info("No questions match your filter criteria.")
    else:
        st.info("Your question bank is empty. Generate some questions to build your collection.")

elif tab == "⚙️ Settings":
    st.markdown('<h1 class="main-title">⚙️ Settings</h1>', unsafe_allow_html=True)
    st.markdown('<p class="subtitle">Configure the application preferences</p>', unsafe_allow_html=True)
    
    with st.expander("API Configuration", expanded=True):
        api_key = st.text_input(
            "Gemini API Key", 
            value=st.session_state.user_settings['api_key'],
            type="password", 
            help="Get your key from Google AI Studio"
        )
        
        model = st.selectbox(
            "Model", 
            ["gemini-1.5-pro-latest", "gemini-1.0-pro"],
            index=0 if st.session_state.user_settings['model'] == "gemini-1.5-pro-latest" else 1,
            help="Select the Gemini model to use for question generation"
        )
    
    with st.expander("Display Settings", expanded=True):
        show_explanations = st.checkbox(
            "Show explanations by default", 
            value=st.session_state.user_settings['show_explanations'],
            help="Always show explanations when displaying questions"
        )
    
    if st.button("Save Settings", use_container_width=True):
        st.session_state.user_settings['api_key'] = api_key
        st.session_state.user_settings['model'] = model
        st.session_state.user_settings['show_explanations'] = show_explanations
        st.success("Settings saved successfully!")
    
    # Add a data management section
    st.markdown("### Data Management")
    
    col1, col2 = st.columns(2)
    
    with col1:
        if st.button("Clear Question Bank", use_container_width=True):
            if 'question_bank' in st.session_state:
                confirm = st.checkbox("Confirm deletion? This cannot be undone.")
                if confirm:
                    st.session_state.question_bank = []
                    st.success("Question bank cleared successfully!")
    
    with col2:
        if st.button("Clear Generation History", use_container_width=True):
            if 'generation_history' in st.session_state:
                confirm = st.checkbox("Confirm deletion? This cannot be undone.", key="confirm_history")
                if confirm:
                    st.session_state.generation_history = []
                    st.success("Generation history cleared successfully!")
    
    # App information
    st.markdown("### About EduGenius")
    st.markdown("""
    EduGenius is an AI-powered application for generating educational questions from various learning materials.
    
    **Version:** 1.0.0  
    **Last Updated:** April 15, 2025  
    **Powered by:** Streamlit and Google Gemini AI
    
    For support or feedback, please contact support@edugenius.ai
    """)
