import streamlit as st
from io import BytesIO
import google.generativeai as genai
import pandas as pd
from datetime import datetime
import matplotlib.pyplot as plt
import seaborn as sns
from wordcloud import WordCloud
from sklearn.feature_extraction.text import TfidfVectorizer
from nltk.tokenize import word_tokenize
from nltk.corpus import stopwords
import nltk
import string

# Download NLTK resources
nltk.download('punkt')
nltk.download('stopwords')

# Set page configuration
st.set_page_config(
    page_title="Lecture2Exam",
    page_icon="🧠",
    layout="wide",
    initial_sidebar_state="collapsed"
)

# Enhanced CSS styling with more vibrant colors and better UI elements
st.markdown("""
<style>
    .main-header {
        font-size: 2.8rem;
        color: #2563EB;
        text-align: center;
        margin-bottom: 2rem;
        font-weight: bold;
        text-shadow: 1px 1px 3px rgba(0,0,0,0.15);
        padding: 1rem;
        background: linear-gradient(90deg, #EFF6FF 0%, #DBEAFE 100%);
        border-radius: 12px;
    }
    .question-box {
        background-color: #F0F9FF;
        border-left: 5px solid #0EA5E9;
        padding: 20px;
        margin-bottom: 25px;
        border-radius: 8px;
        box-shadow: 0 4px 12px rgba(0,0,0,0.08);
        transition: transform 0.2s ease;
    }
    .question-box:hover {
        transform: translateY(-3px);
        box-shadow: 0 6px 16px rgba(0,0,0,0.1);
    }
    .stButton>button {
        background-color: #2563EB;
        color: white;
        font-weight: bold;
        padding: 0.6rem 1.2rem;
        border-radius: 10px;
        transition: all 0.3s;
        border: none;
    }
    .stButton>button:hover {
        background-color: #1D4ED8;
        transform: translateY(-2px);
        box-shadow: 0 4px 12px rgba(0,0,0,0.15);
    }
    .error-message {
        color: #DC2626;
        font-weight: bold;
        padding: 15px;
        background-color: #FEE2E2;
        border-radius: 8px;
        border-left: 5px solid #DC2626;
        margin: 10px 0;
    }
    .subheader {
        font-size: 1.7rem;
        color: #1E40AF;
        margin-top: 1.5rem;
        margin-bottom: 1rem;
        border-bottom: 3px solid #BFDBFE;
        padding-bottom: 0.6rem;
    }
    .question-number {
        font-weight: bold;
        color: #2563EB;
        font-size: 1.25rem;
    }
    .question-text {
        font-size: 1.15rem;
        margin-bottom: 1rem;
        line-height: 1.5;
    }
    .answer {
        margin-top: 0.8rem;
        color: #047857;
        font-weight: bold;
        background-color: #ECFDF5;
        padding: 8px 12px;
        border-radius: 6px;
        display: inline-block;
    }
    .explanation {
        background-color: #F8FAFC;
        padding: 15px;
        border-radius: 8px;
        margin-top: 0.8rem;
        margin-bottom: 1.5rem;
        font-style: italic;
        border: 1px solid #E2E8F0;
    }
    .stats-card {
        background: linear-gradient(145deg, #DBEAFE 0%, #E0F2FE 100%);
        padding: 20px;
        border-radius: 12px;
        margin-bottom: 20px;
        box-shadow: 0 4px 6px rgba(0,0,0,0.05);
        transition: transform 0.2s;
    }
    .stats-card:hover {
        transform: translateY(-5px);
        box-shadow: 0 6px 12px rgba(0,0,0,0.1);
    }
    .stats-card h4 {
        color: #1E40AF;
        font-size: 1.1rem;
        margin-bottom: 8px;
    }
    .stats-card h2 {
        color: #2563EB;
        font-size: 2.5rem;
        margin: 0;
    }
    .file-uploader {
        background-color: #F1F5F9;
        padding: 20px;
        border-radius: 10px;
        border: 2px dashed #CBD5E1;
    }
    .config-section {
        background-color: #F8FAFC;
        padding: 20px;
        border-radius: 10px;
        margin-bottom: 20px;
        box-shadow: 0 2px 8px rgba(0,0,0,0.05);
    }
    .action-button {
        margin-top: 8px;
    }
    .footer {
        text-align: center;
        color: #64748B;
        font-size: 0.9rem;
        padding: 15px;
        margin-top: 30px;
        border-top: 1px solid #E2E8F0;
    }
    /* Custom tab styling */
    .stTabs [data-baseweb="tab-list"] {
        gap: 8px;
    }
    .stTabs [data-baseweb="tab"] {
        height: 50px;
        white-space: pre-wrap;
        background-color: #EFF6FF;
        border-radius: 8px 8px 0 0;
        gap: 1px;
        padding-top: 10px;
        padding-bottom: 10px;
    }
    .stTabs [aria-selected="true"] {
        background-color: #DBEAFE;
        border-bottom: 3px solid #3B82F6;
    }
    /* Graph container styling */
    .graph-container {
        background-color: white;
        padding: 20px;
        border-radius: 10px;
        margin-bottom: 20px;
        box-shadow: 0 2px 8px rgba(0,0,0,0.1);
    }
    /* TF-IDF styling */
    .tfidf-container {
        background-color: #F0FDF4;
        padding: 20px;
        border-radius: 10px;
        margin-bottom: 20px;
        border-left: 5px solid #10B981;
    }
    .tfidf-header {
        color: #065F46;
        font-weight: bold;
        margin-bottom: 15px;
    }
    /* Word cloud styling */
    .wordcloud-container {
        background-color: #F5F3FF;
        padding: 20px;
        border-radius: 10px;
        margin-bottom: 20px;
        border-left: 5px solid #8B5CF6;
    }
</style>
""", unsafe_allow_html=True)

# History tracking
if 'generation_history' not in st.session_state:
    st.session_state.generation_history = []

if 'processed_texts' not in st.session_state:
    st.session_state.processed_texts = []

def extract_text(file):
    try:
        if file.type == "text/plain":
            return file.getvalue().decode("utf-8")
        
        elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            from docx import Document
            doc = Document(BytesIO(file.read()))
            return "\n".join([para.text for para in doc.paragraphs])
        
        elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            from pptx import Presentation
            prs = Presentation(BytesIO(file.read()))
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
        
        elif file.type == "application/pdf":
            import PyPDF2
            pdf_reader = PyPDF2.PdfReader(BytesIO(file.read()))
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
            return text
        
        else:
            st.error("Unsupported file type")
            return ""
    
    except Exception as e:
        st.error(f"Error processing file: {str(e)}")
        return ""

def preprocess_text(text):
    """Tokenize and clean text for TF-IDF analysis"""
    # Tokenize
    tokens = word_tokenize(text.lower())
    
    # Remove punctuation and stopwords
    stop_words = set(stopwords.words('english'))
    tokens = [word for word in tokens if word not in stop_words and word not in string.punctuation]
    
    # Remove numbers and single characters
    tokens = [word for word in tokens if word.isalpha() and len(word) > 1]
    
    return " ".join(tokens)

def calculate_tfidf(texts):
    """Calculate TF-IDF scores for a list of documents"""
    # Preprocess all texts
    processed_texts = [preprocess_text(text) for text in texts]
    
    # Create TF-IDF vectorizer
    vectorizer = TfidfVectorizer(max_features=50)
    tfidf_matrix = vectorizer.fit_transform(processed_texts)
    
    # Get feature names and scores
    feature_names = vectorizer.get_feature_names_out()
    tfidf_scores = tfidf_matrix.toarray()
    
    return feature_names, tfidf_scores

def plot_tfidf(feature_names, tfidf_scores, title="TF-IDF Scores"):
    """Create a bar plot of top TF-IDF terms"""
    fig, ax = plt.subplots(figsize=(10, 6))
    
    # Get average scores across documents
    avg_scores = tfidf_scores.mean(axis=0)
    
    # Create DataFrame for visualization
    df = pd.DataFrame({
        'Term': feature_names,
        'TF-IDF Score': avg_scores
    }).sort_values('TF-IDF Score', ascending=False).head(20)
    
    # Create bar plot
    sns.barplot(x='TF-IDF Score', y='Term', data=df, palette='viridis', ax=ax)
    ax.set_title(title, fontsize=14)
    ax.set_xlabel('TF-IDF Score')
    ax.set_ylabel('')
    
    return fig

def generate_wordcloud(text):
    """Generate a word cloud from text"""
    # Preprocess text
    processed_text = preprocess_text(text)
    
    # Generate word cloud
    wordcloud = WordCloud(width=800, height=400, 
                         background_color='white',
                         colormap='viridis',
                         max_words=50).generate(processed_text)
    
    # Plot
    fig, ax = plt.subplots(figsize=(12, 6))
    ax.imshow(wordcloud, interpolation='bilinear')
    ax.axis('off')
    ax.set_title('Key Terms Word Cloud', fontsize=14)
    
    return fig

def format_questions(raw_questions):
    """Format the questions with better spacing and styling"""
    if not raw_questions:
        return ""
    
    # Split by numbered questions
    import re
    questions = re.split(r'\n\s*(\d+)\.\s+', raw_questions)
    
    if len(questions) <= 1:
        return raw_questions
    
    formatted = ""
    
    # Skip the first empty element if it exists
    start_idx = 1 if questions[0].strip() == "" else 0
    
    for i in range(start_idx, len(questions), 2):
        if i+1 < len(questions):
            question_num = questions[i]
            question_content = questions[i+1].strip()
            
            # Add extra spacing between parts
            question_content = question_content.replace("Question:", "\nQuestion:")
            question_content = question_content.replace("Options:", "\nOptions:")
            question_content = question_content.replace("Correct Answer:", "\n\nCorrect Answer:")
            question_content = question_content.replace("Explanation:", "\n\nExplanation:")
            
            formatted += f'<div class="question-box">\n'
            formatted += f'<span class="question-number">🔍 {question_num}.</span>\n'
            formatted += f'{question_content}\n'
            formatted += f'</div>\n\n'
    
    return formatted

def generate_with_gemini(text, question_type, difficulty, api_key, num_questions=5):
    try:
        genai.configure(api_key=api_key)
        
        # Use the correct model name - gemini-1.5-pro-latest or gemini-1.0-pro
        model = genai.GenerativeModel('gemini-1.5-pro-latest')
        
        prompt = f"""
        Generate {num_questions} {difficulty.lower()} {question_type} questions based on this content.
        Format should be:
        
        1. Question: [question text]
           Options (if MCQ): A) [option1] B) [option2] C) [option3] D) [option4]
           Correct Answer: [correct answer]
           Explanation: [brief explanation]
        
        Content:
        {text[:12000]}  # Using first 12k chars to avoid token limits
        """
        
        response = model.generate_content(prompt)
        
        if not response.text:
            raise ValueError("Empty response from API")
            
        return response.text
    
    except Exception as e:
        st.markdown(f'<div class="error-message">⚠️ API error: {str(e)}</div>', unsafe_allow_html=True)
        st.info("🔑 Please ensure you're using the correct API key and model name")
        st.info("💡 Try using 'gemini-1.0-pro' if this model isn't available")
        return None

def main():
    st.markdown('<h1 class="main-header">🧠 Lecture2Exam ✨</h1>', unsafe_allow_html=True)
    
    # Create tabs for main functionality and history
    tab1, tab2 = st.tabs(["📝 Generate Questions", "📊 History & Analytics"])
    
    with tab1:
        col1, col2 = st.columns([1, 2])
        
        with col1:
            st.markdown('<div class="config-section">', unsafe_allow_html=True)
            st.markdown("### 🔑 API Configuration")
            
            api_key = st.text_input(" API Key", type="password", 
                                   help="Get your key from Google AI Studio")
            st.markdown('</div>', unsafe_allow_html=True)
            
            st.markdown('<div class="file-uploader">', unsafe_allow_html=True)
            st.markdown("### 📄 Upload Learning Material")
            uploaded_file = st.file_uploader("Choose file", 
                                           type=["txt", "docx", "pptx", "pdf"])
            st.markdown('</div>', unsafe_allow_html=True)
            
            if uploaded_file:
                file_info = st.success(f"📑 File uploaded: {uploaded_file.name}")
                
                st.markdown('<div class="config-section">', unsafe_allow_html=True)
                st.markdown("### ⚙️ Question Settings")
                
                question_type = st.selectbox("Question Type 📋", 
                                           ["MCQ", "Fill-in-the-Blank", "Short Answer"])
                
                difficulty_icons = {"Easy": "🟢", "Medium": "🟡", "Hard": "🔴"}
                difficulties = [f"{v} {k}" for k, v in difficulty_icons.items()]
                difficulty_display = st.selectbox("Difficulty Level", difficulties)
                difficulty = difficulty_display.split(' ')[1]  # Extract the actual difficulty
                
                num_questions = st.slider("Number of Questions 🔢", 
                                        min_value=3, max_value=15, value=5,
                                        help="More questions will take longer to generate")
                st.markdown('</div>', unsafe_allow_html=True)
                
                st.markdown('<div class="action-button">', unsafe_allow_html=True)
                gen_col1, gen_col2 = st.columns(2)
                with gen_col1:
                    generate_btn = st.button("✨ Generate Questions", use_container_width=True)
                with gen_col2:
                    clear_btn = st.button("🗑️ Clear Results", use_container_width=True)
                st.markdown('</div>', unsafe_allow_html=True)
                
                if clear_btn:
                    if 'questions' in st.session_state:
                        del st.session_state.questions
                
                if generate_btn:
                    if not api_key:
                        st.error("⚠️ Please enter your API key")
                    else:
                        with st.spinner("🧙‍♂️ Generating intelligent questions..."):
                            text = extract_text(uploaded_file)
                            if text:
                                # Store processed text for TF-IDF analysis
                                st.session_state.processed_texts.append(text)
                                
                                # Try with latest model first, fallback to 1.0 if needed
                                questions = generate_with_gemini(text, question_type, difficulty, api_key, num_questions)
                                if not questions:
                                    # Fallback to gemini-1.0-pro if latest model fails
                                    genai.configure(api_key=api_key)
                                    model = genai.GenerativeModel('gemini-1.0-pro')
                                    questions = generate_with_gemini(text, question_type, difficulty, api_key, num_questions)
                                
                                if questions:
                                    st.session_state.questions = questions
                                    st.session_state.formatted_questions = format_questions(questions)
                                    
                                    # Add to history
                                    st.session_state.generation_history.append({
                                        'timestamp': datetime.now().strftime("%Y-%m-%d %H:%M"),
                                        'file': uploaded_file.name,
                                        'type': question_type,
                                        'difficulty': difficulty,
                                        'count': num_questions,
                                        'questions': questions
                                    })
                                else:
                                    st.session_state.questions = None
        
        with col2:
            if uploaded_file and 'questions' in st.session_state and st.session_state.questions:
                st.markdown(f"<h2 class='subheader'>✅ Generated {question_type} Questions</h2>", unsafe_allow_html=True)
                
                # Show formatted questions with improved spacing
                if 'formatted_questions' in st.session_state:
                    st.markdown(st.session_state.formatted_questions, unsafe_allow_html=True)
                else:
                    st.markdown(st.session_state.questions)
                
                col_a, col_b = st.columns(2)
                with col_a:
                    st.download_button(
                        label="📥 Download Questions",
                        data=BytesIO(st.session_state.questions.encode('utf-8')),
                        file_name=f"{question_type.lower().replace(' ', '_')}_questions.txt",
                        mime="text/plain",
                        use_container_width=True
                    )
                with col_b:
                    if st.button("🔄 Regenerate Questions", use_container_width=True):
                        with st.spinner("🔄 Regenerating questions..."):
                            text = extract_text(uploaded_file)
                            if text:
                                questions = generate_with_gemini(text, question_type, difficulty, api_key, num_questions)
                                if questions:
                                    st.session_state.questions = questions
                                    st.session_state.formatted_questions = format_questions(questions)
            else:
                st.markdown("""
                <div style="background-color: #F0F9FF; padding: 30px; border-radius: 10px; text-align: center; margin-top: 50px;">
                    <h3>🚀 Ready to Create Questions?</h3>
                    <p>Upload your learning material and configure settings to generate AI-powered questions.</p>
                    <p style="font-size: 3rem; margin: 20px 0;">📚 ➡️ 🧠 ➡️ 📝</p>
                    <p>Perfect for teachers, students, and learning professionals.</p>
                </div>
                """, unsafe_allow_html=True)
    
    with tab2:
        st.markdown("<h2 class='subheader'>📊 Generation History & Analytics</h2>", unsafe_allow_html=True)
        
        if not st.session_state.generation_history:
            st.info("📭 No question generation history yet. Generate some questions to see your history.")
        else:
            history_df = pd.DataFrame(st.session_state.generation_history)
            
            # Show stats
            st.markdown("<h3 class='subheader'>📈 Usage Statistics</h3>", unsafe_allow_html=True)
            col1, col2, col3 = st.columns(3)
            with col1:
                st.markdown(f"""
                <div class="stats-card">
                <h4>🔄 Total Generations</h4>
                <h2>{len(history_df)}</h2>
                </div>
                """, unsafe_allow_html=True)
            with col2:
                st.markdown(f"""
                <div class="stats-card">
                <h4>📋 Question Types</h4>
                <h2>{len(history_df['type'].unique())}</h2>
                </div>
                """, unsafe_allow_html=True)
            with col3:
                st.markdown(f"""
                <div class="stats-card">
                <h4>📑 Files Processed</h4>
                <h2>{len(history_df['file'].unique())}</h2>
                </div>
                """, unsafe_allow_html=True)
            
            # Additional stats row
            if len(history_df) > 0:
                col4, col5, col6 = st.columns(3)
                with col4:
                    total_questions = history_df['count'].sum()
                    st.markdown(f"""
                    <div class="stats-card">
                    <h4>🧩 Total Questions Created</h4>
                    <h2>{total_questions}</h2>
                    </div>
                    """, unsafe_allow_html=True)
                with col5:
                    most_common_type = history_df['type'].mode()[0]
                    st.markdown(f"""
                    <div class="stats-card">
                    <h4>⭐ Most Used Type</h4>
                    <h2>{most_common_type}</h2>
                    </div>
                    """, unsafe_allow_html=True)
                with col6:
                    most_common_difficulty = history_df['difficulty'].mode()[0]
                    st.markdown(f"""
                    <div class="stats-card">
                    <h4>🎯 Favorite Difficulty</h4>
                    <h2>{most_common_difficulty}</h2>
                    </div>
                    """, unsafe_allow_html=True)
            
            # Visualizations section
            st.markdown("<h3 class='subheader'>📊 Visual Analytics</h3>", unsafe_allow_html=True)
            
            # Create tabs for different visualizations
            viz_tab1, viz_tab2, viz_tab3 = st.tabs(["📊 Question Types", "📈 Difficulty Levels", "📌 TF-IDF Analysis"])
            
            with viz_tab1:
                st.markdown('<div class="graph-container">', unsafe_allow_html=True)
                fig, ax = plt.subplots(figsize=(10, 5))
                sns.countplot(data=history_df, y='type', order=history_df['type'].value_counts().index, 
                             palette='Blues_r', ax=ax)
                ax.set_title('Question Types Distribution', fontsize=14)
                ax.set_xlabel('Count')
                ax.set_ylabel('Question Type')
                st.pyplot(fig)
                st.markdown('</div>', unsafe_allow_html=True)
            
            with viz_tab2:
                st.markdown('<div class="graph-container">', unsafe_allow_html=True)
                fig, ax = plt.subplots(figsize=(10, 5))
                sns.countplot(data=history_df, y='difficulty', order=['Easy', 'Medium', 'Hard'],
                              palette=['#10B981', '#F59E0B', '#EF4444'], ax=ax)
                ax.set_title('Difficulty Level Distribution', fontsize=14)
                ax.set_xlabel('Count')
                ax.set_ylabel('Difficulty Level')
                st.pyplot(fig)
                st.markdown('</div>', unsafe_allow_html=True)
            
            with viz_tab3:
                if len(st.session_state.processed_texts) > 0:
                    st.markdown('<div class="tfidf-container">', unsafe_allow_html=True)
                    st.markdown('<h4 class="tfidf-header">🔍 Top Important Terms (TF-IDF Analysis)</h4>', unsafe_allow_html=True)
                    
                    # Calculate TF-IDF
                    feature_names, tfidf_scores = calculate_tfidf(st.session_state.processed_texts)
                    
                    # Plot TF-IDF
                    fig = plot_tfidf(feature_names, tfidf_scores)
                    st.pyplot(fig)
                    
                    # Show word cloud
                    st.markdown('<div class="wordcloud-container">', unsafe_allow_html=True)
                    st.markdown('<h4 class="tfidf-header">☁️ Key Terms Word Cloud</h4>', unsafe_allow_html=True)
                    combined_text = " ".join(st.session_state.processed_texts)
                    wordcloud_fig = generate_wordcloud(combined_text)
                    st.pyplot(wordcloud_fig)
                    st.markdown('</div>', unsafe_allow_html=True)
                    
                    st.markdown('</div>', unsafe_allow_html=True)
                else:
                    st.info("No text data available for TF-IDF analysis. Generate questions first.")
            
            # Display history table
            st.markdown("<h3 class='subheader'>📜 Recent Generation Activity</h3>", unsafe_allow_html=True)
            
            # Create a more readable history table
            display_df = history_df[['timestamp', 'file', 'type', 'difficulty', 'count']].copy()
            display_df.columns = ['⏰ Timestamp', '📄 File', '📋 Type', '🎯 Difficulty', '🔢 Count']
            st.dataframe(display_df.tail(10), use_container_width=True)
            
            # Option to clear history
            if st.button("🗑️ Clear History"):
                st.session_state.generation_history = []
                st.session_state.processed_texts = []
                st.experimental_rerun()

    # Footer
    st.markdown("---")
    st.markdown(
        """
        <div class="footer">
            🧠 Lecture2Exam - AI-Powered Assessment Generator<br>
            Make learning more effective with AI-powered assessments ✨<br>
            Made with teamwork of Shreyas, Shaurya and Mahati 🎯
        </div>
        """, 
        unsafe_allow_html=True
    )

if __name__ == "__main__":
    main()
