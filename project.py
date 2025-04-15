import streamlit as st
from io import BytesIO
import google.generativeai as genai
import pandas as pd
from datetime import datetime

# Set page configuration
st.set_page_config(
    page_title="AI Question Generator",
    page_icon="📚",
    layout="wide"
)

# Enhanced CSS styling
st.markdown("""
<style>
    .main-header {
        font-size: 2.5rem;
        color: #1E3A8A;
        text-align: center;
        margin-bottom: 2rem;
        font-weight: bold;
        text-shadow: 1px 1px 2px rgba(0,0,0,0.1);
    }
    .question-box {
        background-color: #EFF6FF;
        border-left: 5px solid #3B82F6;
        padding: 20px;
        margin-bottom: 25px;
        border-radius: 8px;
        box-shadow: 0 4px 6px rgba(0,0,0,0.05);
    }
    .stButton>button {
        background-color: #2563EB;
        color: white;
        font-weight: bold;
        padding: 0.5rem 1rem;
        border-radius: 8px;
        transition: all 0.3s;
    }
    .stButton>button:hover {
        background-color: #1D4ED8;
        transform: translateY(-2px);
        box-shadow: 0 4px 8px rgba(0,0,0,0.1);
    }
    .error-message {
        color: #ff4b4b;
        font-weight: bold;
        padding: 10px;
        background-color: #FECACA;
        border-radius: 5px;
    }
    .subheader {
        font-size: 1.5rem;
        color: #1E3A8A;
        margin-top: 1.5rem;
        margin-bottom: 1rem;
        border-bottom: 2px solid #BFDBFE;
        padding-bottom: 0.5rem;
    }
    .question-number {
        font-weight: bold;
        color: #2563EB;
    }
    .question-text {
        font-size: 1.1rem;
        margin-bottom: 0.8rem;
    }
    .answer {
        margin-top: 0.5rem;
        color: #047857;
        font-weight: bold;
    }
    .explanation {
        background-color: #F3F4F6;
        padding: 10px;
        border-radius: 5px;
        margin-top: 0.5rem;
        margin-bottom: 1.5rem;
        font-style: italic;
    }
    .stats-card {
        background-color: #DBEAFE;
        padding: 15px;
        border-radius: 8px;
        margin-bottom: 15px;
    }
</style>
""", unsafe_allow_html=True)

# History tracking
if 'generation_history' not in st.session_state:
    st.session_state.generation_history = []

def extract_text(file):
    try:
        if file.type == "text/plain":
            return file.getvalue().decode("utf-8")
        
        elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            from docx import Document
            doc = Document(BytesIO(file.read()))
            return "\n".join([para.text for para in doc.paragraphs])
        
        elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            from pptx import Presentation
            prs = Presentation(BytesIO(file.read()))
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
        
        elif file.type == "application/pdf":
            import PyPDF2
            pdf_reader = PyPDF2.PdfReader(BytesIO(file.read()))
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
            return text
        
        else:
            st.error("Unsupported file type")
            return ""
    
    except Exception as e:
        st.error(f"Error processing file: {str(e)}")
        return ""

def format_questions(raw_questions):
    """Format the questions with better spacing and styling"""
    if not raw_questions:
        return ""
    
    # Split by numbered questions
    import re
    questions = re.split(r'\n\s*(\d+)\.\s+', raw_questions)
    
    if len(questions) <= 1:
        return raw_questions
    
    formatted = ""
    
    # Skip the first empty element if it exists
    start_idx = 1 if questions[0].strip() == "" else 0
    
    for i in range(start_idx, len(questions), 2):
        if i+1 < len(questions):
            question_num = questions[i]
            question_content = questions[i+1].strip()
            
            # Add extra spacing between parts
            question_content = question_content.replace("Question:", "\nQuestion:")
            question_content = question_content.replace("Options:", "\nOptions:")
            question_content = question_content.replace("Correct Answer:", "\n\nCorrect Answer:")
            question_content = question_content.replace("Explanation:", "\n\nExplanation:")
            
            formatted += f'<div class="question-box">\n'
            formatted += f'<span class="question-number">{question_num}.</span>\n'
            formatted += f'{question_content}\n'
            formatted += f'</div>\n\n'
    
    return formatted

def generate_with_gemini(text, question_type, difficulty, api_key, num_questions=5):
    try:
        genai.configure(api_key=api_key)
        
        # Use the correct model name - gemini-1.5-pro-latest or gemini-1.0-pro
        model = genai.GenerativeModel('gemini-1.5-pro-latest')
        
        prompt = f"""
        Generate {num_questions} {difficulty.lower()} {question_type} questions based on this content.
        Format should be:
        
        1. Question: [question text]
           Options (if MCQ): A) [option1] B) [option2] C) [option3] D) [option4]
           Correct Answer: [correct answer]
           Explanation: [brief explanation]
        
        Content:
        {text[:12000]}  # Using first 12k chars to avoid token limits
        """
        
        response = model.generate_content(prompt)
        
        if not response.text:
            raise ValueError("Empty response from Gemini API")
            
        return response.text
    
    except Exception as e:
        st.markdown(f'<div class="error-message">Gemini API error: {str(e)}</div>', unsafe_allow_html=True)
        st.info("Please ensure you're using the correct API key and model name")
        st.info("Try using 'gemini-1.0-pro' if this model isn't available")
        return None

def main():
    st.markdown('<h1 class="main-header">📚 AI Question Generator</h1>', unsafe_allow_html=True)
    
    # Create tabs for main functionality and history
    tab1, tab2 = st.tabs(["📝 Generate Questions", "📊 History"])
    
    with tab1:
        col1, col2 = st.columns([1, 2])
        
        with col1:
            st.markdown("### Configuration")
            
            api_key = st.text_input("Gemini API Key", type="password", 
                                   help="Get your key from Google AI Studio")
            
            st.markdown("### Upload Material")
            uploaded_file = st.file_uploader("Choose file", 
                                           type=["txt", "docx", "pptx", "pdf"])
            
            if uploaded_file:
                file_info = st.info(f"File: {uploaded_file.name} ({uploaded_file.type})")
                
                question_type = st.selectbox("Question Type", 
                                           ["MCQ", "Fill-in-the-Blank", "Short Answer"])
                difficulty = st.selectbox("Difficulty", ["Easy", "Medium", "Hard"])
                num_questions = st.slider("Number of Questions", min_value=3, max_value=15, value=5)
                
                gen_col1, gen_col2 = st.columns(2)
                with gen_col1:
                    generate_btn = st.button("✨ Generate Questions", use_container_width=True)
                with gen_col2:
                    clear_btn = st.button("🗑️ Clear", use_container_width=True)
                
                if clear_btn:
                    if 'questions' in st.session_state:
                        del st.session_state.questions
                
                if generate_btn:
                    if not api_key:
                        st.error("Please enter your Gemini API key")
                    else:
                        with st.spinner("Generating questions..."):
                            text = extract_text(uploaded_file)
                            if text:
                                # Try with latest model first, fallback to 1.0 if needed
                                questions = generate_with_gemini(text, question_type, difficulty, api_key, num_questions)
                                if not questions:
                                    # Fallback to gemini-1.0-pro if latest model fails
                                    genai.configure(api_key=api_key)
                                    model = genai.GenerativeModel('gemini-1.0-pro')
                                    questions = generate_with_gemini(text, question_type, difficulty, api_key, num_questions)
                                
                                if questions:
                                    st.session_state.questions = questions
                                    st.session_state.formatted_questions = format_questions(questions)
                                    
                                    # Add to history
                                    st.session_state.generation_history.append({
                                        'timestamp': datetime.now().strftime("%Y-%m-%d %H:%M"),
                                        'file': uploaded_file.name,
                                        'type': question_type,
                                        'difficulty': difficulty,
                                        'count': num_questions,
                                        'questions': questions
                                    })
                                else:
                                    st.session_state.questions = None
        
        with col2:
            if uploaded_file and 'questions' in st.session_state and st.session_state.questions:
                st.markdown(f"<h2 class='subheader'>Generated {question_type} Questions</h2>", unsafe_allow_html=True)
                
                # Show formatted questions with improved spacing
                if 'formatted_questions' in st.session_state:
                    st.markdown(st.session_state.formatted_questions, unsafe_allow_html=True)
                else:
                    st.markdown(st.session_state.questions)
                
                col_a, col_b = st.columns(2)
                with col_a:
                    st.download_button(
                        label="📥 Download Questions",
                        data=BytesIO(st.session_state.questions.encode('utf-8')),
                        file_name=f"{question_type.lower().replace(' ', '_')}_questions.txt",
                        mime="text/plain",
                        use_container_width=True
                    )
                with col_b:
                    if st.button("🔄 Regenerate Questions", use_container_width=True):
                        with st.spinner("Regenerating questions..."):
                            text = extract_text(uploaded_file)
                            if text:
                                questions = generate_with_gemini(text, question_type, difficulty, api_key, num_questions)
                                if questions:
                                    st.session_state.questions = questions
                                    st.session_state.formatted_questions = format_questions(questions)
    
    with tab2:
        st.markdown("<h2 class='subheader'>Generation History</h2>", unsafe_allow_html=True)
        
        if not st.session_state.generation_history:
            st.info("No question generation history yet. Generate some questions to see your history.")
        else:
            history_df = pd.DataFrame(st.session_state.generation_history)
            
            # Show stats
            st.markdown("<h3 class='subheader'>Generation Stats</h3>", unsafe_allow_html=True)
            col1, col2, col3 = st.columns(3)
            with col1:
                st.markdown(f"""
                <div class="stats-card">
                <h4>Total Generations</h4>
                <h2>{len(history_df)}</h2>
                </div>
                """, unsafe_allow_html=True)
            with col2:
                st.markdown(f"""
                <div class="stats-card">
                <h4>Question Types</h4>
                <h2>{len(history_df['type'].unique())}</h2>
                </div>
                """, unsafe_allow_html=True)
            with col3:
                st.markdown(f"""
                <div class="stats-card">
                <h4>Files Processed</h4>
                <h2>{len(history_df['file'].unique())}</h2>
                </div>
                """, unsafe_allow_html=True)
            
            # Display history table
            st.markdown("<h3 class='subheader'>Recent Generations</h3>", unsafe_allow_html=True)
            
            # Create a more readable history table
            display_df = history_df[['timestamp', 'file', 'type', 'difficulty', 'count']].copy()
            display_df.columns = ['Timestamp', 'File', 'Question Type', 'Difficulty', '# Questions']
            st.dataframe(display_df.tail(10), use_container_width=True)
            
            # Option to clear history
            if st.button("🗑️ Clear History"):
                st.session_state.generation_history = []
                st.experimental_rerun()

    # Footer
    st.markdown("---")
    st.markdown(
        """
        <div style="text-align: center; color: #6B7280; font-size: 0.8rem;">
            AI Question Generator | Powered by Gemini | Created with Streamlit
        </div>
        """, 
        unsafe_allow_html=True
    )

if __name__ == "__main__":
    main()
